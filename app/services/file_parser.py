"""Parse individual files into content the AI can consume.

Supports: code, PDF (always converted to images), DOCX, images, Jupyter notebooks.
ALL PDFs are converted to images for consistent vision-based analysis.
"""
from __future__ import annotations

import base64
import json
import logging
import mimetypes
import re
import zipfile
from pathlib import Path
from typing import Optional
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

CODE_EXTENSIONS: set[str] = {
    ".py", ".java", ".cpp", ".c", ".js", ".ts", ".cs", ".go", ".rb",
    ".php", ".swift", ".kt", ".scala", ".rs", ".sh", ".sql", ".html",
    ".css", ".r", ".m", ".h", ".jsx", ".tsx", ".vue", ".svelte",
    ".dart", ".lua", ".perl", ".pl", ".ps1", ".bat", ".asm", ".swift",
}

IMAGE_EXTENSIONS: set[str] = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg"
}

TEXT_EXTENSIONS: set[str] = {
    ".txt", ".md", ".csv", ".xml", ".json", ".rtf", ".log", ".yaml", ".yml"
}

DOCUMENT_EXTENSIONS: set[str] = {
    ".doc", ".docx", ".odt", ".pdf", ".xls", ".xlsx", ".ppt", ".pptx"
}

ARCHIVE_EXTENSIONS: set[str] = {
    ".zip", ".rar", ".7z", ".tar", ".gz", ".tgz", ".bz2"
}

TRANSIENT_FILE_PREFIXES: tuple[str, ...] = ("~$", "._")
TRANSIENT_FILE_NAMES: set[str] = {".ds_store", "thumbs.db", "desktop.ini"}


def _is_transient_or_system_file(file_path: Path) -> bool:
    """Return True for Office lock files and OS metadata files."""
    name = file_path.name
    lower_name = name.lower()

    if name.startswith(TRANSIENT_FILE_PREFIXES):
        return True

    return lower_name in TRANSIENT_FILE_NAMES


def parse_file(file_path: Path) -> dict:
    """Return a dict describing the file content for the AI.

    Keys:
        filename, type, content (text or base64), size, error (if any)
    """
    if isinstance(file_path, str):
        file_path = Path(file_path)
    
    ext = file_path.suffix.lower()
    result: dict = {
        "filename": file_path.name,
        "extension": ext,
        "size": file_path.stat().st_size if file_path.exists() else 0,
    }

    if _is_transient_or_system_file(file_path):
        result.update({
            "type": "skipped",
            "content": None,
            "note": "Transient/system file skipped"
        })
        return result

    try:
        if ext in CODE_EXTENSIONS:
            result.update(_parse_code(file_path))
        elif ext == ".pdf":
            result.update(_parse_pdf_to_images(file_path))
        elif ext == ".docx":
            result.update(_parse_docx(file_path))
        elif ext == ".doc":
            result.update(_parse_doc(file_path))
        elif ext in IMAGE_EXTENSIONS:
            result.update(_parse_image(file_path))
        elif ext == ".ipynb":
            result.update(_parse_notebook(file_path))
        elif ext in TEXT_EXTENSIONS:
            result.update(_parse_text(file_path))
        elif ext in ARCHIVE_EXTENSIONS:
            result.update({"type": "archive", "content": None, "note": "Archive file - contents extracted separately"})
        elif ext in [".xlsx", ".xls"]:
            result.update(_parse_excel(file_path))
        elif ext in [".pptx", ".ppt"]:
            result.update(_parse_powerpoint(file_path))
        else:
            result.update(_parse_as_text_fallback(file_path))
    except Exception as exc:
        logger.exception("Error parsing %s", file_path.name)
        result["type"] = "error"
        result["content"] = None
        result["error"] = str(exc)

    return result


def _parse_code(file_path: Path) -> dict:
    text = file_path.read_text(encoding="utf-8", errors="replace")
    return {"type": "code", "content": text, "language": _detect_language(file_path.suffix)}


def _parse_text(file_path: Path) -> dict:
    text = file_path.read_text(encoding="utf-8", errors="replace")
    return {"type": "text", "content": text}


def _parse_as_text_fallback(file_path: Path) -> dict:
    """Try to parse unknown file types as text."""
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
        if text.strip():
            return {"type": "text", "content": text, "note": "Parsed as text (unknown extension)"}
    except Exception:
        pass
    return {"type": "binary", "content": None, "note": "Binary file - cannot extract text"}


def _parse_pdf_to_images(file_path: Path, dpi: int = 200, max_pages: int = 20) -> dict:
    """Convert ALL PDF pages to base64 PNG images for vision input.
    
    This ensures consistent handling regardless of whether the PDF contains:
    - Text only
    - Scanned documents
    - Images/diagrams
    - Handwritten content
    - Mixed content
    
    Args:
        file_path: Path to the PDF file
        dpi: Resolution for rendering (higher = better quality but larger)
        max_pages: Maximum number of pages to convert (to avoid huge payloads)
    
    Returns:
        dict with type='pdf_images' and content as list of base64 image strings
    """
    import fitz  # PyMuPDF
    
    doc = fitz.open(str(file_path))
    images: list[str] = []
    page_count = len(doc)
    
    logger.info(f"Converting PDF '{file_path.name}' ({page_count} pages) to images for vision analysis")
    
    for page_num in range(min(page_count, max_pages)):
        page = doc[page_num]
        pix = page.get_pixmap(dpi=dpi)
        img_bytes = pix.tobytes("png")
        b64 = base64.b64encode(img_bytes).decode()
        images.append(b64)
    
    if page_count > max_pages:
        logger.warning(f"PDF '{file_path.name}' has {page_count} pages, only converted first {max_pages}")
    
    doc.close()
    
    return {
        "type": "pdf_images",
        "content": images,
        "page_count": min(page_count, max_pages),
        "total_pages": page_count,
        "dpi": dpi,
        "note": f"PDF converted to {len(images)} page images for vision analysis"
    }


def _parse_docx(file_path: Path) -> dict:
    """Parse DOCX - extract both text and convert to images for vision."""
    from docx import Document
    
    doc = Document(str(file_path))
    parts: list[str] = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    
    text_content = "\n".join(parts)
    
    return {
        "type": "docx", 
        "content": text_content,
        "has_images": _docx_has_images(doc),
        "note": "Word document - text extracted"
    }


def _docx_has_images(doc) -> bool:
    """Check if DOCX contains images."""
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                return True
    except Exception:
        pass
    return False


def _parse_doc(file_path: Path) -> dict:
    """Parse old .doc format - try text extraction."""
    try:
        import subprocess
        result = subprocess.run(
            ['textutil', '-convert', 'txt', '-stdout', str(file_path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return {"type": "doc", "content": result.stdout, "note": "Legacy Word document"}
    except Exception:
        pass
    
    return {"type": "doc", "content": None, "error": "Could not parse .doc file - try converting to .docx"}


def _parse_excel(file_path: Path) -> dict:
    """Parse Excel files using openpyxl/xlrd fallback."""
    ext = file_path.suffix.lower()
    errors: list[str] = []

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(filename=str(file_path), data_only=True, read_only=True)
            parts: list[str] = []
            sheet_names: list[str] = []

            for sheet in workbook.worksheets:
                sheet_names.append(sheet.title)
                parts.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(value.strip() for value in values):
                        parts.append(" | ".join(values))

            workbook.close()
            return {
                "type": "excel",
                "content": "\n".join(parts),
                "note": "Excel spreadsheet",
                "sheets": sheet_names,
                "extraction_method": "openpyxl",
            }
        except Exception as exc:
            errors.append(f"openpyxl: {exc}")

    if ext == ".xls":
        try:
            import xlrd

            workbook = xlrd.open_workbook(str(file_path))
            parts: list[str] = []
            sheet_names: list[str] = []

            for sheet in workbook.sheets():
                sheet_names.append(sheet.name)
                parts.append(f"=== Sheet: {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    if any(value.strip() for value in values):
                        parts.append(" | ".join(values))

            return {
                "type": "excel",
                "content": "\n".join(parts),
                "note": "Excel spreadsheet",
                "sheets": sheet_names,
                "extraction_method": "xlrd",
            }
        except Exception as exc:
            errors.append(f"xlrd: {exc}")

    return {"type": "excel", "content": None, "error": "; ".join(errors) if errors else "Unsupported or corrupted Excel file"}


def _parse_powerpoint(file_path: Path) -> dict:
    """Parse PowerPoint files using python-pptx and XML fallback."""
    pptx_error: Optional[Exception] = None

    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return {
            "type": "powerpoint",
            "content": "\n".join(parts),
            "note": "PowerPoint presentation",
            "extraction_method": "python-pptx",
        }
    except Exception as exc:
        pptx_error = exc

    xml_error: Optional[Exception] = None
    if file_path.suffix.lower() == ".pptx":
        try:
            with zipfile.ZipFile(file_path, "r") as archive:
                slide_paths = sorted(
                    [
                        name for name in archive.namelist()
                        if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                    ],
                    key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1))
                    if re.search(r"slide(\d+)\.xml$", name) else 10**9,
                )

                parts: list[str] = []
                for idx, slide_path in enumerate(slide_paths, 1):
                    root = ET.fromstring(archive.read(slide_path))
                    texts = [
                        node.text.strip()
                        for node in root.iter()
                        if node.tag.endswith("}t") and node.text and node.text.strip()
                    ]
                    parts.append(f"=== Slide {idx} ===")
                    if texts:
                        parts.append("\n".join(texts))

                if slide_paths:
                    return {
                        "type": "powerpoint",
                        "content": "\n".join(parts),
                        "note": "PowerPoint presentation",
                        "extraction_method": "pptx_xml_fallback",
                    }
        except Exception as exc:
            xml_error = exc

    errors: list[str] = []
    if pptx_error:
        errors.append(f"python-pptx: {pptx_error}")
    if xml_error:
        errors.append(f"xml_fallback: {xml_error}")

    return {
        "type": "powerpoint",
        "content": None,
        "error": "; ".join(errors) if errors else "Unsupported or corrupted PowerPoint file",
    }


def _parse_image(file_path: Path) -> dict:
    """Base64-encode an image for AI vision input."""
    raw = file_path.read_bytes()
    b64 = base64.b64encode(raw).decode()
    ext = file_path.suffix.lower().lstrip(".")
    if ext == "jpg":
        ext = "jpeg"
    if ext == "tif":
        ext = "tiff"
    return {
        "type": "image", 
        "content": b64, 
        "media_type": f"image/{ext}",
        "size_bytes": len(raw)
    }


def _parse_notebook(file_path: Path) -> dict:
    """Convert .ipynb to structured text with code + markdown cells."""
    import nbformat

    nb = nbformat.read(str(file_path), as_version=4)
    parts: list[str] = []

    for i, cell in enumerate(nb.cells, 1):
        if cell.cell_type == "code":
            parts.append(f"# --- Code Cell {i} ---\n{cell.source}")
            if cell.get("outputs"):
                for output in cell.outputs:
                    if output.get("text"):
                        parts.append(f"# Output:\n{output['text']}")
        elif cell.cell_type == "markdown":
            parts.append(f"# --- Markdown Cell {i} ---\n{cell.source}")

    return {"type": "notebook", "content": "\n\n".join(parts)}


def _detect_language(ext: str) -> str:
    """Detect programming language from extension."""
    lang_map = {
        ".py": "python",
        ".java": "java",
        ".cpp": "cpp",
        ".c": "c",
        ".h": "c",
        ".js": "javascript",
        ".ts": "typescript",
        ".jsx": "jsx",
        ".tsx": "tsx",
        ".cs": "csharp",
        ".go": "go",
        ".rb": "ruby",
        ".php": "php",
        ".swift": "swift",
        ".kt": "kotlin",
        ".scala": "scala",
        ".rs": "rust",
        ".sh": "bash",
        ".sql": "sql",
        ".html": "html",
        ".css": "css",
        ".r": "r",
        ".m": "matlab",
        ".lua": "lua",
        ".pl": "perl",
    }
    return lang_map.get(ext.lower(), "unknown")


def get_file_type_summary(files: list[dict]) -> dict:
    """Get a summary of file types in a submission."""
    summary = {
        "code": [],
        "images": [],
        "pdfs": [],
        "documents": [],
        "text": [],
        "notebooks": [],
        "other": [],
        "errors": [],
        "total_count": len(files),
        "total_size": 0,
    }
    
    for f in files:
        file_type = f.get("type", "unknown")
        filename = f.get("filename", "unknown")
        size = f.get("size", 0)
        summary["total_size"] += size
        
        if file_type == "code":
            summary["code"].append(filename)
        elif file_type == "image":
            summary["images"].append(filename)
        elif file_type == "pdf_images":
            summary["pdfs"].append(f"{filename} ({f.get('page_count', '?')} pages)")
        elif file_type in ("docx", "doc", "excel", "powerpoint"):
            summary["documents"].append(filename)
        elif file_type == "text":
            summary["text"].append(filename)
        elif file_type == "notebook":
            summary["notebooks"].append(filename)
        elif file_type == "error":
            summary["errors"].append(f"{filename}: {f.get('error', 'unknown error')}")
        else:
            summary["other"].append(filename)
    
    return summary
