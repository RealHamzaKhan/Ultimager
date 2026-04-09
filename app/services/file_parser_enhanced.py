"""
Enhanced file ingestion pipeline with full transparency logging.

Handles ALL submission formats with detailed extraction reports:
- Text content → extract and preserve
- Images (standalone, PDF pages, DOCX embedded) → send to LLM as vision inputs
- Mixed files → split into text + image parts

Provides per-student ingestion reports for complete transparency.
"""
from __future__ import annotations

import base64
import io
import json
import logging
import mimetypes
import re
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple, Union
import hashlib
from collections import defaultdict
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

# ============================================================================
# INTELLIGENT CODE CHUNKING - Respects function/class boundaries
# ============================================================================

# Language-specific patterns for detecting code boundaries
CODE_BOUNDARY_PATTERNS: Dict[str, Dict[str, re.Pattern]] = {
    "python": {
        "class": re.compile(r'^class\s+\w+.*?:\s*$', re.MULTILINE),
        "function": re.compile(r'^def\s+\w+.*?:\s*$', re.MULTILINE),
        "async_function": re.compile(r'^async\s+def\s+\w+.*?:\s*$', re.MULTILINE),
        "import": re.compile(r'^(?:from\s+[\w.]+\s+)?import\s+[\w.,\s]+', re.MULTILINE),
    },
    "java": {
        "class": re.compile(r'^(?:public\s+|private\s+|protected\s+)?(?:static\s+)?class\s+\w+', re.MULTILINE),
        "method": re.compile(r'^(?:public\s+|private\s+|protected\s+)?(?:static\s+)?(?:void|int|String|boolean|double|float|long|char|byte|short|Object)\s+\w+\s*\([^)]*\)\s*(?:throws\s+\w+(?:,\s*\w+)*)?\s*\{', re.MULTILINE),
        "import": re.compile(r'^import\s+[\w.*]+;', re.MULTILINE),
    },
    "javascript": {
        "class": re.compile(r'^class\s+\w+.*?\{', re.MULTILINE),
        "function": re.compile(r'^(?:const|let|var|function)\s+\w+\s*=.*?(?:=>|\bfunction\b)', re.MULTILINE),
        "function_decl": re.compile(r'^function\s+\w+\s*\(', re.MULTILINE),
        "import": re.compile(r'^(?:import\s+.*?from\s+[\'"].*?[\'"]|require\s*\(\s*[\'"].*?[\'"]\s*\))', re.MULTILINE),
    },
    "typescript": {
        "class": re.compile(r'^class\s+\w+.*?\{', re.MULTILINE),
        "function": re.compile(r'^(?:const|let|var|function)\s+\w+\s*[:(].*?(?:=>|\bfunction\b)', re.MULTILINE),
        "interface": re.compile(r'^interface\s+\w+.*?\{', re.MULTILINE),
        "import": re.compile(r'^import\s+.*?from\s+[\'"].*?[\'"]', re.MULTILINE),
    },
    "cpp": {
        "class": re.compile(r'^(?:class|struct)\s+\w+.*?\{', re.MULTILINE),
        "function": re.compile(r'^(?:void|int|double|float|char|bool|string|auto)\s+\w+\s*\([^)]*\)\s*(?:const)?\s*\{', re.MULTILINE),
        "include": re.compile(r'^#include\s*<[^>]+>', re.MULTILINE),
    },
    "c": {
        "class": re.compile(r'^(?:struct|union|enum)\s+\w+.*?\{', re.MULTILINE),
        "function": re.compile(r'^(?:void|int|double|float|char|bool|long|short)\s+\w+\s*\([^)]*\)\s*\{', re.MULTILINE),
        "include": re.compile(r'^#include\s*[<"][^>"]+[>"]', re.MULTILINE),
    },
    "go": {
        "function": re.compile(r'^func\s+(?:\(\w+\s+\*?\w+\)\s+)?\w+\s*\(', re.MULTILINE),
        "type": re.compile(r'^type\s+\w+\s+(?:struct|interface)\s*\{', re.MULTILINE),
        "import": re.compile(r'^import\s+(?:\(|[\'"])', re.MULTILINE),
    },
    "rust": {
        "function": re.compile(r'^fn\s+\w+.*?\{', re.MULTILINE),
        "struct": re.compile(r'^struct\s+\w+.*?\{', re.MULTILINE),
        "impl": re.compile(r'^impl(?:\s+\w+)?.*?\{', re.MULTILINE),
        "use": re.compile(r'^use\s+[\w:]+', re.MULTILINE),
    },
    "ruby": {
        "class": re.compile(r'^class\s+\w+.*$', re.MULTILINE),
        "def": re.compile(r'^def\s+\w+.*$', re.MULTILINE),
        "module": re.compile(r'^module\s+\w+.*$', re.MULTILINE),
        "require": re.compile(r'^(?:require|require_relative)\s+[\'"].*?[\'"]', re.MULTILINE),
    },
    "php": {
        "class": re.compile(r'^(?:abstract\s+)?class\s+\w+.*?\{', re.MULTILINE),
        "function": re.compile(r'^(?:public|private|protected|static)\s+function\s+\w+\s*\(', re.MULTILINE),
        "use": re.compile(r'^use\s+[\w\\]+;', re.MULTILINE),
    },
    "swift": {
        "class": re.compile(r'^(?:class|struct|enum|protocol)\s+\w+.*?\{', re.MULTILINE),
        "func": re.compile(r'^func\s+\w+.*?\{', re.MULTILINE),
        "import": re.compile(r'^import\s+(?:struct|class|func|typealias|protocol)?\s*\w+', re.MULTILINE),
    },
    "kotlin": {
        "class": re.compile(r'^(?:class|interface|object)\s+\w+.*?\{', re.MULTILINE),
        "fun": re.compile(r'^fun\s+\w+.*?\{', re.MULTILINE),
        "import": re.compile(r'^import\s+[\w.]+', re.MULTILINE),
    },
}

# Default fallback patterns for unknown languages
DEFAULT_BOUNDARY_PATTERNS = {
    "function": re.compile(r'^func\s+', re.MULTILINE),
    "import": re.compile(r'^(?:import|require|include)\s+', re.MULTILINE),
}


def _get_language_patterns(language: str) -> Dict[str, re.Pattern]:
    """Get code boundary patterns for a specific language."""
    return CODE_BOUNDARY_PATTERNS.get(language.lower(), DEFAULT_BOUNDARY_PATTERNS)


def _smart_chunk_code(text: str, language: str, max_chars: int) -> Tuple[str, Dict[str, Any]]:
    """
    Intelligently chunk code content respecting function/class boundaries.

    Returns (chunked_content, chunk_metadata) where chunk_metadata contains:
    - total_chunks: total number of chunks
    - chunk_boundaries: list of (chunk_id, start_pos, end_pos, code_unit_name)
    - truncated: whether content was truncated
    - truncated_at: position where truncation happened
    - missing_code_units: list of code units that were cut off
    """
    if len(text) <= max_chars:
        return text, {
            "total_chunks": 1,
            "chunk_boundaries": [(0, 0, len(text), "full_file")],
            "truncated": False,
            "truncated_at": None,
            "missing_code_units": [],
            "original_length": len(text),
        }

    patterns = _get_language_patterns(language)

    # Find all code unit positions
    code_units = []

    # Find function/class definitions
    for unit_type, pattern in patterns.items():
        if unit_type == "import":
            continue  # Handle imports separately
        for match in pattern.finditer(text):
            # Extract the name if possible
            name = match.group(0).strip()[:50]
            code_units.append((match.start(), match.end(), unit_type, name))

    # Sort by position
    code_units.sort(key=lambda x: x[0])

    # Find import statements
    import_pattern = patterns.get("import")
    imports = []
    if import_pattern:
        for match in import_pattern.finditer(text):
            imports.append(match.group(0).strip())

    # Chunk content – respect max_chars as a *total output budget*
    chunks = []
    current_pos = 0
    chunk_boundaries = []
    missing_code_units = []
    total_emitted = 0
    separator = "\n### Code continues: "
    truncation_warning_budget = 120  # reserve space for the trailing warning

    while current_pos < len(text):
        # How many chars can we still emit?
        budget = max_chars - total_emitted - truncation_warning_budget
        if budget <= 0:
            break

        if current_pos + budget >= len(text):
            # Fits in remaining space
            chunk_text = text[current_pos:]
            chunks.append(chunk_text)
            chunk_boundaries.append((len(chunks) - 1, current_pos, len(text), "final"))
            total_emitted += len(chunk_text)
            break

        # Find the best break point within budget
        best_break = -1
        best_unit_name = "chunk_end"

        for pos, end_pos, unit_type, name in code_units:
            if current_pos <= pos < current_pos + budget:
                if pos > best_break:
                    best_break = pos
                    best_unit_name = name or f"{unit_type}_boundary"

        if best_break > current_pos:
            chunk_text = text[current_pos:best_break]
            chunks.append(chunk_text)
            chunk_boundaries.append((len(chunks) - 1, current_pos, best_break, best_unit_name))
            total_emitted += len(chunk_text)
            current_pos = best_break
        else:
            # No natural break point, hard break
            chunk_text = text[current_pos:current_pos + budget]
            chunks.append(chunk_text)
            chunk_boundaries.append((len(chunks) - 1, current_pos, current_pos + budget, "hard_break"))
            total_emitted += len(chunk_text)
            current_pos += budget

    # Check for incomplete code units
    last_chunk_end = chunk_boundaries[-1][2] if chunk_boundaries else 0
    for pos, end_pos, unit_type, name in code_units:
        if pos < last_chunk_end and end_pos > last_chunk_end:
            missing_code_units.append({
                "type": unit_type,
                "name": name,
                "started_at": pos,
                "unfinished": True
            })

    # Assemble final content with proper annotations
    final_parts = []
    for i, chunk in enumerate(chunks):
        if i > 0:
            final_parts.append("\n[...CONTINUED FROM PREVIOUS SECTION...]\n")

        # Add structure summary at chunk boundaries
        if chunk_boundaries[i][3] not in ("full_file", "final"):
            final_parts.append(f"\n### Code continues: {chunk_boundaries[i][3]} ###\n")

        final_parts.append(chunk)

    # Add truncation warning if needed
    if len(text) > max_chars:
        final_parts.append(f"\n\n[WARNING: CONTENT TRUNCATED at {max_chars} chars. Original: {len(text)} chars]")
        if missing_code_units:
            unit_names = [u["name"][:30] for u in missing_code_units]
            final_parts.append(f"\n[TRUNCATED UNFINISHED CODE: {', '.join(unit_names)}]")

    result = "".join(final_parts)

    return result, {
        "total_chunks": len(chunks),
        "chunk_boundaries": chunk_boundaries,
        "truncated": len(text) > max_chars,
        "truncated_at": len(text) if len(text) > max_chars else None,
        "missing_code_units": missing_code_units,
        "imports_found": imports,
        "original_length": len(text),
    }


# ============================================================================
# FILE DEPENDENCY DETECTION
# ============================================================================

def detect_file_dependencies(file_contents: List['ExtractedContent']) -> Dict[str, Dict[str, Any]]:
    """
    Detect dependencies between files based on imports, requires, includes, etc.

    Returns a dictionary mapping filenames to their dependencies:
    {
        "main.py": {
            "imports": ["helper.py", "utils.py"],
            "imported_by": ["app.py"],
            "dependency_type": "python_import"
        },
        ...
    }
    """
    dependency_map = {}
    all_files = {fc.filename.lower(): fc for fc in file_contents if fc.filename}

    # Patterns for detecting dependencies
    dependency_patterns = {
        # Python: from X import Y / import X
        "python": [
            re.compile(r'from\s+([\w./]+)\s+import', re.MULTILINE),
            re.compile(r'import\s+([\w.,\s]+)', re.MULTILINE),
        ],
        # JavaScript/TypeScript: import X from 'Y' / require('Y')
        "javascript": [
            re.compile(r"import\s+.*?\s+from\s+['\"]([^'\"]+)['\"]", re.MULTILINE),
            re.compile(r"import\s+['\"]([^'\"]+)['\"]", re.MULTILINE),
            re.compile(r"require\s*\(\s*['\"]([^'\"]+)['\"]\s*\)", re.MULTILINE),
        ],
        # Java: import com.package.Class
        "java": [
            re.compile(r'import\s+([\w.]+);', re.MULTILINE),
        ],
        # C/C++: #include "file.h" / <file.h>
        "cpp": [
            re.compile(r'#include\s*["<]([^">]+)[">]', re.MULTILINE),
        ],
        # Go: import "package"
        "go": [
            re.compile(r'import\s+["\（]([^"\）]+)["\）]', re.MULTILINE),
        ],
        # Ruby: require 'file'
        "ruby": [
            re.compile(r"require\s+['\"]([^'\"]+)['\"]", re.MULTILINE),
            re.compile(r"require_relative\s+['\"]([^'\"]+)['\"]", re.MULTILINE),
        ],
        # PHP: require 'file.php'
        "php": [
            re.compile(r"(?:require|include|require_once|include_once)\s+['\"]([^'\"]+)['\"]", re.MULTILINE),
        ],
    }

    for content in file_contents:
        if not content.text_content or content.file_type not in ("code", "notebook"):
            continue

        lang = content.metadata.get("language", "unknown")
        filename = content.filename.lower()

        dependencies = {
            "imports": [],
            "imported_by": [],
            "dependency_type": None,
            "has_local_imports": False,
            "missing_imports": [],
        }

        # Get patterns for this language
        patterns = dependency_patterns.get(lang.lower(), [])

        # Also check for JS/TS in .js/.ts files
        if not patterns and content.filename:
            ext = content.filename.rsplit('.', 1)[-1] if '.' in content.filename else ''
            if ext in ('js', 'jsx', 'ts', 'tsx', 'mjs'):
                patterns = dependency_patterns.get("javascript", [])

        # Extract dependencies
        for pattern in patterns:
            for match in pattern.finditer(content.text_content):
                matched_text = match.group(1).strip() if match.groups() else match.group(0)
                # Clean up the matched text
                matched_text = re.sub(r'[{}\s,;].*', '', matched_text)  # Remove { }, commas, semicolons

                if matched_text and not matched_text.startswith(('http', 'www', '//', '#')):
                    # Handle relative paths
                    base_name = matched_text.rsplit('/', 1)[-1]  # Get last component

                    # Try to match with actual files
                    for fname in all_files:
                        if (base_name in fname or
                            base_name.replace('.py', '') in fname or
                            base_name.replace('.js', '') in fname or
                            base_name.replace('.ts', '') in fname):
                            if fname != filename and fname not in dependencies["imports"]:
                                dependencies["imports"].append(fname)

                    # Check if it's a local import (not a standard library/package)
                    if matched_text and not matched_text.startswith(('sys', 'os', 're', 'json', 'math', 'typing')):
                        if matched_text not in dependencies["imports"]:
                            # Check if it could be a local file
                            dependencies["has_local_imports"] = True

        # Check for missing imports
        for imp in dependencies["imports"]:
            if imp not in all_files:
                dependencies["missing_imports"].append(imp)
                dependencies["imports"].remove(imp)

        dependency_map[content.filename.lower()] = dependencies

    # Build reverse relationship (imported_by)
    for fname, deps in dependency_map.items():
        for imp in deps.get("imports", []):
            if imp in dependency_map:
                if fname not in dependency_map[imp]["imported_by"]:
                    dependency_map[imp]["imported_by"].append(fname)

    return dependency_map


def detect_image_sequences(file_contents: List['ExtractedContent']) -> List[Dict[str, Any]]:
    """
    Detect image sequences that might form a complete picture (e.g., numbered files).

    Returns list of image sequence groups:
    [
        {
            "base_name": "diagram",
            "extension": ".png",
            "files": ["diagram1.png", "diagram2.png", "diagram3.png"],
            "sequence": [1, 2, 3],
            "total_count": 3
        },
        ...
    ]
    """
    # Group images by base name and extension
    image_files: Dict[Tuple[str, str], List[Tuple[str, int]]] = defaultdict(list)

    for content in file_contents:
        if content.file_type != "image":
            continue

        filename = content.filename
        # Match pattern: nameNNN.ext or name_NNN.ext
        match = re.match(r'^(.+?)[_ ]?(\d+)\.(\w+)$', filename, re.IGNORECASE)
        if match:
            base_name = match.group(1).lower()
            num = int(match.group(2))
            ext = match.group(3).lower()
            image_files[(base_name, ext)].append((filename, num))
        else:
            # Also check for pageNNN pattern (like scanned docs)
            match = re.match(r'^(.+?)page[-_]?(\d+)\.(\w+)$', filename, re.IGNORECASE)
            if match:
                base_name = match.group(1).lower()
                num = int(match.group(2))
                ext = match.group(3).lower()
                image_files[(base_name, ext)].append((filename, num))

    sequences = []
    for (base_name, ext), files in image_files.items():
        if len(files) > 1:  # Only groups with multiple images
            sorted_files = sorted(files, key=lambda x: x[1])
            sequences.append({
                "base_name": base_name,
                "extension": f".{ext}",
                "files": [f[0] for f in sorted_files],
                "sequence": [f[1] for f in sorted_files],
                "total_count": len(files),
                "is_complete": all(i in [f[1] for f in sorted_files] for i in range(1, len(files) + 1))
            })

    return sequences

# Extended file type support
CODE_EXTENSIONS: set[str] = {
    ".py", ".java", ".cpp", ".c", ".js", ".ts", ".cs", ".go", ".rb",
    ".php", ".swift", ".kt", ".scala", ".rs", ".sh", ".sql", ".html",
    ".css", ".r", ".m", ".h", ".jsx", ".tsx", ".vue", ".svelte",
    ".dart", ".lua", ".perl", ".pl", ".ps1", ".bat", ".asm",
}

IMAGE_EXTENSIONS: set[str] = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg"
}

TEXT_EXTENSIONS: set[str] = {
    ".txt", ".md", ".csv", ".xml", ".json", ".rtf", ".log", ".yaml", ".yml",
    ".rst", ".tex", ".cfg", ".ini", ".properties"
}

DOCUMENT_EXTENSIONS: set[str] = {
    ".doc", ".docx", ".odt", ".pdf", ".xls", ".xlsx", ".ppt", ".pptx",
    ".epub", ".pages", ".numbers", ".key"
}

ARCHIVE_EXTENSIONS: set[str] = {
    ".zip", ".rar", ".7z", ".tar", ".gz", ".tgz", ".bz2", ".xz"
}

TRANSIENT_FILE_PREFIXES: tuple[str, ...] = ("~$", "._")
TRANSIENT_FILE_NAMES: set[str] = {".ds_store", "thumbs.db", "desktop.ini"}


def _is_transient_or_system_file(file_path: Path) -> bool:
    """Return True for editor/OS/Office transient files that should never be parsed."""
    name = file_path.name
    lower_name = name.lower()

    if name.startswith(TRANSIENT_FILE_PREFIXES):
        return True

    return lower_name in TRANSIENT_FILE_NAMES

@dataclass
class ExtractedContent:
    """Represents extracted content from a file."""
    filename: str
    file_type: str  # 'code', 'text', 'image', 'pdf', 'docx', 'mixed', 'error'
    text_content: Optional[str] = None
    images: List[Dict[str, Any]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None
    extraction_method: str = "unknown"
    size_bytes: int = 0
    
    def to_dict(self) -> dict:
        return {
            "filename": self.filename,
            "file_type": self.file_type,
            "text_content_preview": self.text_content[:200] + "..." if self.text_content and len(self.text_content) > 200 else self.text_content,
            "text_length": len(self.text_content) if self.text_content else 0,
            "image_count": len(self.images),
            "metadata": self.metadata,
            "error": self.error,
            "extraction_method": self.extraction_method,
            "size_bytes": self.size_bytes,
        }


@dataclass
class IngestionReport:
    """Complete ingestion report for a student submission."""
    student_id: str
    session_id: int
    timestamp: str
    files_received: List[Dict[str, Any]] = field(default_factory=list)
    files_parsed: List[Dict[str, Any]] = field(default_factory=list)
    files_failed: List[Dict[str, Any]] = field(default_factory=list)
    total_text_chars: int = 0
    total_images: int = 0
    content_truncated: bool = False
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    
    def to_dict(self) -> dict:
        return {
            "student_id": self.student_id,
            "session_id": self.session_id,
            "timestamp": self.timestamp,
            "files_received": self.files_received,
            "files_parsed": self.files_parsed,
            "files_failed": self.files_failed,
            "total_text_chars": self.total_text_chars,
            "total_images": self.total_images,
            "content_truncated": self.content_truncated,
            "errors": self.errors,
            "warnings": self.warnings,
            "summary": {
                "received": len(self.files_received),
                "parsed": len(self.files_parsed),
                "failed": len(self.files_failed),
            }
        }


# ============================================================================
# D7 FIX: FILE MAGIC BYTE VALIDATION
# ============================================================================
# Detects when file extension doesn't match actual content type.

_MAGIC_SIGNATURES: Dict[bytes, str] = {
    b'%PDF':                  '.pdf',
    b'\x50\x4b\x03\x04':     '.zip',   # ZIP (also docx/xlsx/pptx)
    b'\x89PNG\r\n\x1a\n':    '.png',
    b'\xff\xd8\xff':          '.jpg',
    b'GIF87a':                '.gif',
    b'GIF89a':                '.gif',
    b'RIFF':                  '.webp',  # needs WEBP check at offset 8
    b'BM':                    '.bmp',
    b'PK\x03\x04':           '.zip',
}

# ZIP-based formats identified by internal file patterns
_ZIP_SUBTYPES = {
    'word/document.xml': '.docx',
    'xl/workbook.xml': '.xlsx',
    'ppt/presentation.xml': '.pptx',
}


def _validate_magic_bytes(file_path: Path, claimed_ext: str) -> Optional[str]:
    """Check if the file's magic bytes match its claimed extension.
    Returns the detected extension if it differs, or None if it matches or is undetectable.
    """
    try:
        with open(file_path, 'rb') as f:
            header = f.read(16)
    except Exception:
        return None

    if len(header) < 4:
        return None

    detected = None
    for sig, ext in _MAGIC_SIGNATURES.items():
        if header[:len(sig)] == sig:
            detected = ext
            break

    if detected is None:
        return None

    # If it's a ZIP, probe internal structure for docx/xlsx/pptx
    if detected == '.zip' and claimed_ext not in ('.zip',):
        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                names = zf.namelist()
                for internal_path, sub_ext in _ZIP_SUBTYPES.items():
                    if internal_path in names:
                        detected = sub_ext
                        break
        except Exception:
            pass

    # Return None if it matches the claimed extension (no mismatch)
    if detected == claimed_ext:
        return None
    # ZIP-based formats: .docx is a .zip, so don't flag docx as zip
    if detected == '.zip' and claimed_ext in ('.docx', '.xlsx', '.pptx', '.jar', '.ipynb'):
        return None
    # Image subtypes: jpg/jpeg are equivalent
    if detected == '.jpg' and claimed_ext in ('.jpg', '.jpeg'):
        return None
    if detected == '.jpeg' and claimed_ext in ('.jpg', '.jpeg'):
        return None

    return detected


def parse_file_with_report(file_path: Path, max_text_chars: int = 80000,
                          max_image_size: int = 5*1024*1024) -> Tuple[ExtractedContent, IngestionReport]:
    """
    Parse a single file with full transparency logging.
    
    Returns tuple of (ExtractedContent, mini_report_for_this_file)
    """
    if isinstance(file_path, str):
        file_path = Path(file_path)
    
    filename = file_path.name
    ext = file_path.suffix.lower()
    size = file_path.stat().st_size if file_path.exists() else 0
    
    report = {
        "filename": filename,
        "size_bytes": size,
        "extension": ext,
        "status": "pending"
    }
    
    try:
        if not file_path.exists():
            report["status"] = "failed"
            report["error"] = "File not found"
            return ExtractedContent(
                filename=filename,
                file_type="error",
                error="File not found",
                size_bytes=0
            ), report

        if _is_transient_or_system_file(file_path):
            report["status"] = "skipped"
            report["warning"] = "Transient/system file skipped"
            return ExtractedContent(
                filename=filename,
                file_type="skipped",
                extraction_method="skipped",
                size_bytes=size,
                metadata={"reason": "transient_system_file"},
            ), report

        # D7 FIX: Magic byte validation — detect extension mismatch
        actual_type = _validate_magic_bytes(file_path, ext)
        if actual_type and actual_type != ext:
            logger.warning(
                f"Magic byte mismatch: {filename} claims {ext} but looks like {actual_type}"
            )
            report["warning"] = f"Extension mismatch: file appears to be {actual_type}"
            ext = actual_type  # Use the actual detected type for routing

        # Route to appropriate parser
        if ext in CODE_EXTENSIONS:
            content = _parse_code_file(file_path, max_text_chars)
            report["status"] = "parsed"
            report["extraction_method"] = "text_utf8"
            report["text_length"] = len(content.text_content) if content.text_content else 0
            return content, report
            
        elif ext in TEXT_EXTENSIONS:
            content = _parse_text_file(file_path, max_text_chars)
            report["status"] = "parsed"
            report["extraction_method"] = "text_utf8"
            report["text_length"] = len(content.text_content) if content.text_content else 0
            return content, report
            
        elif ext == ".pdf":
            content = _parse_pdf_mixed(file_path, max_text_chars)
            report["status"] = "parsed"
            report["extraction_method"] = content.extraction_method
            report["text_length"] = len(content.text_content) if content.text_content else 0
            report["image_count"] = len(content.images)
            return content, report
            
        elif ext == ".docx":
            content = _parse_docx_mixed(file_path, max_text_chars)
            report["status"] = "parsed"
            report["extraction_method"] = content.extraction_method
            report["text_length"] = len(content.text_content) if content.text_content else 0
            report["image_count"] = len(content.images)
            return content, report
            
        elif ext in IMAGE_EXTENSIONS:
            content = _parse_image_file(file_path, max_image_size)
            report["status"] = "parsed"
            report["extraction_method"] = "image_base64"
            report["image_count"] = 1
            return content, report
            
        elif ext == ".ipynb":
            content = _parse_notebook_file(file_path, max_text_chars)
            report["status"] = "parsed"
            report["extraction_method"] = "nbconvert"
            report["text_length"] = len(content.text_content) if content.text_content else 0
            return content, report
            
        elif ext in [".xlsx", ".xls"]:
            content = _parse_excel_file(file_path)
            if content.file_type == "error":
                report["status"] = "failed"
                report["error"] = content.error or "Excel parsing failed"
            else:
                report["status"] = "parsed"
                report["extraction_method"] = content.extraction_method
                report["text_length"] = len(content.text_content) if content.text_content else 0
            return content, report
            
        elif ext in [".pptx", ".ppt"]:
            content = _parse_powerpoint_file(file_path)
            if content.file_type == "error":
                report["status"] = "failed"
                report["error"] = content.error or "PowerPoint parsing failed"
            else:
                report["status"] = "parsed"
                report["extraction_method"] = content.extraction_method
                report["text_length"] = len(content.text_content) if content.text_content else 0
            return content, report
            
        elif ext in ARCHIVE_EXTENSIONS:
            report["status"] = "skipped"
            report["warning"] = "Archive files should be extracted before processing"
            return ExtractedContent(
                filename=filename,
                file_type="archive",
                text_content=None,
                extraction_method="skipped",
                size_bytes=size,
                metadata={"note": "Archive file - contents extracted separately"}
            ), report
            
        else:
            # Try as text fallback
            content = _parse_as_text_fallback(file_path, max_text_chars)
            if content.text_content:
                report["status"] = "parsed"
                report["extraction_method"] = "text_fallback"
                report["text_length"] = len(content.text_content)
                return content, report
            else:
                report["status"] = "skipped"
                report["warning"] = f"Unsupported file type: {ext}"
                return ExtractedContent(
                    filename=filename,
                    file_type="unsupported",
                    extraction_method="skipped",
                    size_bytes=size,
                    metadata={"extension": ext}
                ), report
                
    except Exception as e:
        logger.exception(f"Error parsing {filename}")
        report["status"] = "failed"
        report["error"] = str(e)
        return ExtractedContent(
            filename=filename,
            file_type="error",
            error=str(e),
            size_bytes=size
        ), report


def _parse_code_file(file_path: Path, max_chars: int) -> ExtractedContent:
    """Parse code files with intelligent chunking that respects function/class boundaries."""
    text = file_path.read_text(encoding="utf-8", errors="replace")
    language = _detect_language(file_path.suffix)

    # Use smart chunking to respect function/class boundaries instead of naive truncation
    chunked_text, chunk_meta = _smart_chunk_code(text, language, max_chars)

    return ExtractedContent(
        filename=file_path.name,
        file_type="code",
        text_content=chunked_text,
        extraction_method="smart_chunk" if chunk_meta.get("truncated") else "text_utf8",
        size_bytes=file_path.stat().st_size,
        metadata={
            "language": language,
            "truncated": chunk_meta.get("truncated", False),
            "original_length": chunk_meta.get("original_length", len(text)),
            "total_chunks": chunk_meta.get("total_chunks", 1),
            "missing_code_units": chunk_meta.get("missing_code_units", []),
        }
    )


def _parse_text_file(file_path: Path, max_chars: int) -> ExtractedContent:
    """Parse plain text files."""
    text = file_path.read_text(encoding="utf-8", errors="replace")
    truncated = len(text) > max_chars
    
    return ExtractedContent(
        filename=file_path.name,
        file_type="text",
        text_content=text[:max_chars] if truncated else text,
        extraction_method="text_utf8",
        size_bytes=file_path.stat().st_size,
        metadata={
            "truncated": truncated,
            "original_length": len(text)
        }
    )


def _parse_as_text_fallback(file_path: Path, max_chars: int) -> ExtractedContent:
    """Try to parse unknown files as text."""
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
        if text.strip():
            truncated = len(text) > max_chars
            return ExtractedContent(
                filename=file_path.name,
                file_type="text",
                text_content=text[:max_chars] if truncated else text,
                extraction_method="text_fallback",
                size_bytes=file_path.stat().st_size,
                metadata={
                    "note": "Parsed as text (unknown extension)",
                    "truncated": truncated,
                    "original_length": len(text)
                }
            )
    except Exception:
        pass
    
    return ExtractedContent(
        filename=file_path.name,
        file_type="binary",
        extraction_method="failed",
        size_bytes=file_path.stat().st_size,
        metadata={"note": "Binary file - cannot extract text"}
    )


def _compute_image_hash(img_bytes: bytes) -> int:
    """Compute a perceptual average hash (64-bit) for image deduplication."""
    try:
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(img_bytes)).convert("L").resize((8, 8), PILImage.Resampling.LANCZOS)
        pixels = list(img.getdata())
        avg = sum(pixels) / len(pixels)
        return sum(1 << i for i, p in enumerate(pixels) if p >= avg)
    except Exception:
        return int(hashlib.md5(img_bytes[:4096]).hexdigest()[:16], 16)


def _is_duplicate_image(new_hash: int, existing_hashes: list, threshold: int = 5) -> bool:
    """Check if image is perceptually similar to any existing image (Hamming distance)."""
    for h in existing_hashes:
        hamming = bin(new_hash ^ h).count("1")
        if hamming <= threshold:
            return True
    return False


def _parse_pdf_mixed(file_path: Path, max_text_chars: int = 80000,
                     dpi: int = 200, max_pages: int = 60) -> ExtractedContent:
    """
    Parse PDF with mixed content strategy:
    1. Extract text from all pages
    2. Convert pages to images for vision analysis
    3. Extract embedded high-res images (diagrams, figures) from PDF XObjects
    4. Deduplicate embedded images vs page renders using perceptual hashing
    5. Return both text and all unique images
    """
    import fitz  # PyMuPDF

    doc = fitz.open(str(file_path))
    page_count = len(doc)

    # Extract text
    text_parts = []
    for page_num in range(min(page_count, max_pages)):
        page = doc[page_num]
        text = page.get_text()
        if text.strip():
            text_parts.append(f"\n--- Page {page_num + 1} ---\n{text}")

    full_text = "\n".join(text_parts)
    text_truncated = len(full_text) > max_text_chars

    images = []
    image_hashes: list[int] = []
    MAX_IMAGE_RENDER_PAGES = 30  # Cap page-to-image rendering to limit memory usage
    pages_converted = min(page_count, max_pages)
    pages_rendered = min(pages_converted, MAX_IMAGE_RENDER_PAGES)
    embedded_count = 0
    deduplicated_count = 0

    # Phase 1: Render pages to images (capped to avoid memory blowup on large PDFs)
    for page_num in range(pages_rendered):
        try:
            page = doc[page_num]
            pix = page.get_pixmap(dpi=dpi)
            img_bytes = pix.tobytes("png")
            b64 = base64.b64encode(img_bytes).decode()
            img_hash = _compute_image_hash(img_bytes)
            image_hashes.append(img_hash)

            images.append({
                "page": page_num + 1,
                "base64": b64,
                "media_type": "image/png",
                "size_bytes": len(img_bytes),
                "source": "page_render",
            })
        except Exception as e:
            logger.warning(f"Failed to convert PDF page {page_num + 1}: {e}")

    # Phase 2: Extract embedded XObject images (diagrams, figures, photos at native res)
    for page_num in range(pages_converted):
        try:
            page = doc[page_num]
            image_list = page.get_images(full=True)

            for img_info in image_list:
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue

                    img_bytes = base_image["image"]
                    img_ext = base_image.get("ext", "png")
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)

                    # Skip tiny images (icons, bullets, decorations)
                    if width < 50 or height < 50:
                        continue

                    # Resize very large embedded images (> 5MB)
                    if len(img_bytes) > 5 * 1024 * 1024:
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(io.BytesIO(img_bytes))
                            pil_img.thumbnail((2048, 2048), PILImage.Resampling.LANCZOS)
                            buf = io.BytesIO()
                            pil_img.save(buf, format="PNG")
                            img_bytes = buf.getvalue()
                            img_ext = "png"
                        except Exception:
                            continue

                    # Deduplicate against page renders and previously seen embedded images
                    img_hash = _compute_image_hash(img_bytes)
                    if _is_duplicate_image(img_hash, image_hashes):
                        deduplicated_count += 1
                        continue

                    image_hashes.append(img_hash)
                    media_type = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"
                    b64 = base64.b64encode(img_bytes).decode()
                    embedded_count += 1

                    images.append({
                        "page": page_num + 1,
                        "base64": b64,
                        "media_type": media_type,
                        "size_bytes": len(img_bytes),
                        "source": "embedded_xobject",
                        "native_resolution": f"{width}x{height}",
                        "embedded_id": embedded_count,
                    })
                except Exception as e:
                    logger.debug(f"Failed to extract embedded image xref={xref} p{page_num+1}: {e}")
        except Exception as e:
            logger.debug(f"Failed to scan page {page_num+1} for embedded images: {e}")

    doc.close()

    has_text = len(full_text.strip()) > 50
    has_images = len(images) > 0

    if has_text and has_images:
        extraction_method = "mixed_text_and_vision"
    elif has_images:
        extraction_method = "vision_only"
    else:
        extraction_method = "text_only"

    return ExtractedContent(
        filename=file_path.name,
        file_type="pdf",
        text_content=full_text[:max_text_chars] if text_truncated else full_text,
        images=images,
        extraction_method=extraction_method,
        size_bytes=file_path.stat().st_size,
        metadata={
            "total_pages": page_count,
            "pages_converted": pages_converted,
            "pages_exceeded": page_count > max_pages,
            "text_extracted": has_text,
            "images_extracted": has_images,
            "page_render_images": pages_converted,
            "embedded_images_extracted": embedded_count,
            "deduplicated_images": deduplicated_count,
            "total_images": len(images),
            "truncated": text_truncated,
            "original_text_length": len(full_text),
        }
    )


def _parse_docx_mixed(file_path: Path, max_text_chars: int = 80000) -> ExtractedContent:
    """
    Parse DOCX with mixed content strategy:
    1. Extract text from paragraphs and tables
    2. Extract embedded images
    3. Return both
    """
    from docx import Document
    from docx.oxml.ns import nsmap
    
    doc = Document(str(file_path))
    
    # Extract text
    text_parts = []
    
    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    
    # Tables
    for table_idx, table in enumerate(doc.tables):
        text_parts.append(f"\n--- Table {table_idx + 1} ---")
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text_parts.append(row_text)
    
    full_text = "\n".join(text_parts)
    text_truncated = len(full_text) > max_text_chars
    
    # Extract embedded images
    images = []
    try:
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    
                    # Determine image format
                    content_type = image_part.content_type
                    if "png" in content_type:
                        media_type = "image/png"
                        ext = "png"
                    elif "jpeg" in content_type or "jpg" in content_type:
                        media_type = "image/jpeg"
                        ext = "jpg"
                    else:
                        media_type = "image/png"
                        ext = "png"
                    
                    b64 = base64.b64encode(image_bytes).decode()
                    image_count += 1
                    
                    images.append({
                        "embedded_id": image_count,
                        "base64": b64,
                        "media_type": media_type,
                        "size_bytes": len(image_bytes),
                        "description": f"Embedded image {image_count}"
                    })
                except Exception as e:
                    logger.warning(f"Failed to extract embedded image: {e}")
    except Exception as e:
        logger.warning(f"Failed to process DOCX images: {e}")
    
    # Determine extraction method
    has_text = len(full_text.strip()) > 0
    has_images = len(images) > 0
    
    if has_text and has_images:
        extraction_method = "mixed_text_and_images"
    elif has_images:
        extraction_method = "images_only"
    else:
        extraction_method = "text_only"
    
    return ExtractedContent(
        filename=file_path.name,
        file_type="docx",
        text_content=full_text[:max_text_chars] if text_truncated else full_text,
        images=images,
        extraction_method=extraction_method,
        size_bytes=file_path.stat().st_size,
        metadata={
            "paragraphs": len([p for p in doc.paragraphs if p.text.strip()]),
            "tables": len(doc.tables),
            "embedded_images": len(images),
            "truncated": text_truncated,
            "original_text_length": len(full_text)
        }
    )


def _parse_image_file(file_path: Path, max_size: int = 5*1024*1024) -> ExtractedContent:
    """Parse image files for vision input.

    Always caps resolution at 1536px on the longest side to keep API payloads
    manageable while retaining enough detail for handwriting and diagrams.
    """
    from PIL import Image

    size = file_path.stat().st_size
    MAX_DIM = 1536  # Max pixels on longest side — enough for vision models

    ext = file_path.suffix.lower().lstrip(".")
    if ext == "jpg":
        ext = "jpeg"
    if ext == "tif":
        ext = "tiff"

    # SVG cannot be processed by PIL — read raw bytes and treat as PNG for downstream
    PIL_UNSUPPORTED_FORMATS = {"svg"}
    if ext in PIL_UNSUPPORTED_FORMATS:
        img_bytes = file_path.read_bytes()
        b64 = base64.b64encode(img_bytes).decode()
        return ExtractedContent(
            filename=file_path.name,
            file_type="image",
            images=[{
                "base64": b64,
                "media_type": f"image/svg+xml",
                "size_bytes": len(img_bytes),
                "original_size": file_path.stat().st_size,
                "source": "raw_file",
            }],
            extraction_method="raw_bytes",
            size_bytes=file_path.stat().st_size,
            metadata={"format": "svg", "note": "SVG not processed by PIL"},
        )

    resized = False
    try:
        with Image.open(file_path) as img:
            orig_width, orig_height = img.size
            # Always resize if either dimension exceeds MAX_DIM or file is too large
            needs_resize = (
                max(orig_width, orig_height) > MAX_DIM
                or size > max_size
            )
            if needs_resize:
                img.thumbnail((MAX_DIM, MAX_DIM), Image.Resampling.LANCZOS)
                resized = True

            buffer = io.BytesIO()
            img_format = "JPEG" if ext in ("jpeg", "jpg") else ext.upper()
            save_kwargs = {}
            if img_format == "JPEG":
                save_kwargs["quality"] = 85  # Good quality, smaller size
                save_kwargs["optimize"] = True
            img.save(buffer, format=img_format, **save_kwargs)
            img_bytes = buffer.getvalue()
            final_width, final_height = img.size
    except Exception as e:
        # Fallback: read raw bytes if PIL fails
        img_bytes = file_path.read_bytes()
        orig_width, orig_height = 0, 0
        final_width, final_height = 0, 0

    b64 = base64.b64encode(img_bytes).decode()

    return ExtractedContent(
        filename=file_path.name,
        file_type="image",
        images=[{
            "base64": b64,
            "media_type": f"image/{ext}",
            "size_bytes": len(img_bytes),
            "original_size": size,
            "resized": resized,
            "dimensions": {"width": final_width, "height": final_height},
            "original_dimensions": {"width": orig_width, "height": orig_height},
        }],
        extraction_method="image_base64",
        size_bytes=size,
        metadata={
            "dimensions": {"width": final_width, "height": final_height},
            "original_dimensions": {"width": orig_width, "height": orig_height},
            "resized": resized
        }
    )


def _parse_notebook_file(file_path: Path, max_chars: int) -> ExtractedContent:
    """Parse Jupyter notebooks — extracts code, markdown, text outputs, AND output images."""
    import nbformat

    nb = nbformat.read(str(file_path), as_version=4)
    parts = []
    images = []
    output_image_count = 0

    for i, cell in enumerate(nb.cells, 1):
        if cell.cell_type == "code":
            parts.append(f"# --- Code Cell {i} ---\n{cell.source}")
            if cell.get("outputs"):
                for out_idx, output in enumerate(cell.outputs):
                    # Text output (stdout, stderr, plain text)
                    if output.get("text"):
                        parts.append(f"# Output:\n{output['text']}")
                    elif output.get("output_type") == "stream":
                        text = output.get("text", "")
                        if text:
                            parts.append(f"# Output ({output.get('name', 'stdout')}):\n{text}")

                    # Image outputs (matplotlib plots, display_data, execute_result)
                    data = output.get("data", {})
                    if not data and output.get("output_type") not in ("display_data", "execute_result"):
                        continue

                    # Check for PNG image (most common: matplotlib, seaborn, plotly static)
                    if "image/png" in data:
                        b64_data = data["image/png"]
                        # nbformat stores base64 directly (no prefix)
                        if isinstance(b64_data, str):
                            b64_clean = b64_data.replace("\n", "").strip()
                            try:
                                img_bytes = base64.b64decode(b64_clean)
                                output_image_count += 1
                                images.append({
                                    "base64": b64_clean,
                                    "media_type": "image/png",
                                    "size_bytes": len(img_bytes),
                                    "source": f"notebook_cell_{i}_output_{out_idx}",
                                    "description": f"Output image from code cell {i}",
                                })
                            except Exception as e:
                                logger.debug(f"Failed to decode notebook PNG cell {i}: {e}")

                    # Check for JPEG image
                    elif "image/jpeg" in data:
                        b64_data = data["image/jpeg"]
                        if isinstance(b64_data, str):
                            b64_clean = b64_data.replace("\n", "").strip()
                            try:
                                img_bytes = base64.b64decode(b64_clean)
                                output_image_count += 1
                                images.append({
                                    "base64": b64_clean,
                                    "media_type": "image/jpeg",
                                    "size_bytes": len(img_bytes),
                                    "source": f"notebook_cell_{i}_output_{out_idx}",
                                    "description": f"Output image from code cell {i}",
                                })
                            except Exception as e:
                                logger.debug(f"Failed to decode notebook JPEG cell {i}: {e}")

                    # Check for SVG (render to PNG if cairosvg available)
                    elif "image/svg+xml" in data:
                        svg_data = data["image/svg+xml"]
                        if isinstance(svg_data, str):
                            try:
                                import cairosvg
                                png_bytes = cairosvg.svg2png(bytestring=svg_data.encode("utf-8"))
                                b64_png = base64.b64encode(png_bytes).decode()
                                output_image_count += 1
                                images.append({
                                    "base64": b64_png,
                                    "media_type": "image/png",
                                    "size_bytes": len(png_bytes),
                                    "source": f"notebook_cell_{i}_output_{out_idx}_svg",
                                    "description": f"SVG output rendered from code cell {i}",
                                })
                            except ImportError:
                                logger.debug("cairosvg not installed, skipping SVG render")
                            except Exception as e:
                                logger.debug(f"Failed to render SVG from cell {i}: {e}")

                    # Include text/plain representation as fallback
                    if "text/plain" in data and "image/png" not in data and "image/jpeg" not in data:
                        plain = data["text/plain"]
                        if isinstance(plain, str) and plain.strip():
                            parts.append(f"# Output:\n{plain}")

        elif cell.cell_type == "markdown":
            parts.append(f"# --- Markdown Cell {i} ---\n{cell.source}")

    full_text = "\n\n".join(parts)
    truncated = len(full_text) > max_chars

    return ExtractedContent(
        filename=file_path.name,
        file_type="notebook",
        text_content=full_text[:max_chars] if truncated else full_text,
        images=images if images else [],
        extraction_method="nbformat_with_images" if images else "nbformat",
        size_bytes=file_path.stat().st_size,
        metadata={
            "cells": len(nb.cells),
            "code_cells": sum(1 for c in nb.cells if c.cell_type == "code"),
            "output_images_extracted": output_image_count,
            "truncated": truncated,
            "original_length": len(full_text),
        }
    )


def _parse_excel_file(file_path: Path) -> ExtractedContent:
    """Parse Excel files — extracts cell values, AND embedded images/charts from xlsx."""
    ext = file_path.suffix.lower()
    errors: list[str] = []

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook

            # First pass: read cell data (read_only mode for speed)
            workbook_ro = load_workbook(filename=str(file_path), data_only=True, read_only=True)
            parts: list[str] = []
            sheet_names: list[str] = []

            for sheet in workbook_ro.worksheets:
                sheet_names.append(sheet.title)
                parts.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(value.strip() for value in values):
                        parts.append(" | ".join(values))

            workbook_ro.close()

            # Second pass: extract embedded images (requires non-read-only mode)
            images: list[dict] = []
            chart_count = 0
            try:
                workbook_full = load_workbook(filename=str(file_path))
                for sheet in workbook_full.worksheets:
                    # Extract embedded images
                    for img_obj in getattr(sheet, '_images', []):
                        try:
                            img_data = img_obj._data()
                            if len(img_data) < 500:
                                continue
                            b64 = base64.b64encode(img_data).decode()
                            images.append({
                                "base64": b64,
                                "media_type": "image/png",
                                "size_bytes": len(img_data),
                                "source": f"excel_sheet_{sheet.title}_image",
                            })
                        except Exception as e:
                            logger.debug(f"Failed to extract Excel image: {e}")

                    # Count charts (for metadata — chart rendering requires matplotlib bridge)
                    chart_count += len(getattr(sheet, '_charts', []))

                workbook_full.close()
            except Exception as e:
                logger.debug(f"Non-critical: could not extract Excel images: {e}")

            # Also extract images from xl/media/ in the xlsx ZIP archive
            if not images:
                try:
                    with zipfile.ZipFile(file_path, "r") as zf:
                        media_files = [n for n in zf.namelist() if n.startswith("xl/media/")]
                        for mf in media_files:
                            mf_ext = Path(mf).suffix.lower()
                            if mf_ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".emf", ".wmf"):
                                try:
                                    img_bytes = zf.read(mf)
                                    if len(img_bytes) < 500:
                                        continue

                                    # Handle EMF/WMF (Windows Metafiles) - skip if can't convert
                                    if mf_ext in (".emf", ".wmf"):
                                        continue

                                    mt = f"image/{mf_ext.lstrip('.')}"
                                    if mf_ext == ".jpg":
                                        mt = "image/jpeg"
                                    b64 = base64.b64encode(img_bytes).decode()
                                    images.append({
                                        "base64": b64,
                                        "media_type": mt,
                                        "size_bytes": len(img_bytes),
                                        "source": f"excel_media_{Path(mf).name}",
                                    })
                                except Exception:
                                    pass
                except Exception:
                    pass

            return ExtractedContent(
                filename=file_path.name,
                file_type="excel",
                text_content="\n".join(parts),
                images=images if images else [],
                extraction_method="openpyxl" + ("_with_images" if images else ""),
                size_bytes=file_path.stat().st_size,
                metadata={
                    "sheets": sheet_names,
                    "embedded_images": len(images),
                    "charts_detected": chart_count,
                },
            )
        except Exception as exc:
            errors.append(f"openpyxl: {exc}")

    if ext == ".xls":
        try:
            import xlrd

            workbook = xlrd.open_workbook(str(file_path))
            parts_xls: list[str] = []
            sheet_names_xls: list[str] = []

            for sheet in workbook.sheets():
                sheet_names_xls.append(sheet.name)
                parts_xls.append(f"=== Sheet: {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    if any(value.strip() for value in values):
                        parts_xls.append(" | ".join(values))

            return ExtractedContent(
                filename=file_path.name,
                file_type="excel",
                text_content="\n".join(parts_xls),
                extraction_method="xlrd",
                size_bytes=file_path.stat().st_size,
                metadata={"sheets": sheet_names_xls},
            )
        except Exception as exc:
            errors.append(f"xlrd: {exc}")

    return ExtractedContent(
        filename=file_path.name,
        file_type="error",
        error="; ".join(errors) if errors else "Unsupported or corrupted Excel file",
        extraction_method="failed",
        size_bytes=file_path.stat().st_size,
    )


def _parse_powerpoint_file(file_path: Path) -> ExtractedContent:
    """Parse PowerPoint files — extracts text, speaker notes, AND embedded images."""
    pptx_error: Optional[Exception] = None
    images: list[dict] = []

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(str(file_path))
        parts = []
        image_count = 0

        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)

                # Extract embedded images from picture shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_part = shape.image
                        img_bytes = image_part.blob
                        content_type = image_part.content_type or "image/png"

                        # Skip tiny images (icons)
                        if len(img_bytes) < 1000:
                            continue

                        # Resize if > 5MB
                        if len(img_bytes) > 5 * 1024 * 1024:
                            try:
                                from PIL import Image as PILImage
                                pil_img = PILImage.open(io.BytesIO(img_bytes))
                                pil_img.thumbnail((2048, 2048), PILImage.Resampling.LANCZOS)
                                buf = io.BytesIO()
                                pil_img.save(buf, format="PNG")
                                img_bytes = buf.getvalue()
                                content_type = "image/png"
                            except Exception:
                                continue

                        b64 = base64.b64encode(img_bytes).decode()
                        image_count += 1
                        images.append({
                            "base64": b64,
                            "media_type": content_type,
                            "size_bytes": len(img_bytes),
                            "source": f"slide_{i}_picture",
                            "description": f"Image from slide {i}",
                        })
                    except Exception as e:
                        logger.debug(f"Failed to extract PPTX image from slide {i}: {e}")

                # Extract images from group shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for grp_shape in shape.shapes:
                            if hasattr(grp_shape, 'image'):
                                img_bytes = grp_shape.image.blob
                                if len(img_bytes) < 1000:
                                    continue
                                b64 = base64.b64encode(img_bytes).decode()
                                ct = grp_shape.image.content_type or "image/png"
                                image_count += 1
                                images.append({
                                    "base64": b64,
                                    "media_type": ct,
                                    "size_bytes": len(img_bytes),
                                    "source": f"slide_{i}_group_picture",
                                })
                    except Exception:
                        pass

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Speaker Notes]: {notes_text}")

        return ExtractedContent(
            filename=file_path.name,
            file_type="powerpoint",
            text_content="\n".join(parts),
            images=images if images else [],
            extraction_method="python-pptx" + ("_with_images" if images else ""),
            size_bytes=file_path.stat().st_size,
            metadata={
                "slides": len(prs.slides),
                "embedded_images": image_count,
            }
        )
    except Exception as exc:
        pptx_error = exc

    # XML fallback (text only, no images)
    xml_error: Optional[Exception] = None
    if file_path.suffix.lower() == ".pptx":
        try:
            with zipfile.ZipFile(file_path, "r") as archive:
                slide_paths = sorted(
                    [
                        name for name in archive.namelist()
                        if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                    ],
                    key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1))
                    if re.search(r"slide(\d+)\.xml$", name) else 10**9,
                )

                parts_fb: list[str] = []
                fb_images: list[dict] = []

                for idx, slide_path in enumerate(slide_paths, 1):
                    xml_data = archive.read(slide_path)
                    root = ET.fromstring(xml_data)
                    texts = [
                        node.text.strip()
                        for node in root.iter()
                        if node.tag.endswith("}t") and node.text and node.text.strip()
                    ]
                    parts_fb.append(f"=== Slide {idx} ===")
                    if texts:
                        parts_fb.append("\n".join(texts))

                # Also try extracting images from the ppt/media/ folder
                media_files = [n for n in archive.namelist() if n.startswith("ppt/media/")]
                for mf in media_files:
                    ext = Path(mf).suffix.lower()
                    if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
                        try:
                            img_bytes = archive.read(mf)
                            if len(img_bytes) < 1000:
                                continue
                            media_type = f"image/{ext.lstrip('.')}"
                            if ext == ".jpg":
                                media_type = "image/jpeg"
                            b64 = base64.b64encode(img_bytes).decode()
                            fb_images.append({
                                "base64": b64,
                                "media_type": media_type,
                                "size_bytes": len(img_bytes),
                                "source": f"pptx_media_{Path(mf).name}",
                            })
                        except Exception:
                            pass

                if slide_paths:
                    return ExtractedContent(
                        filename=file_path.name,
                        file_type="powerpoint",
                        text_content="\n".join(parts_fb),
                        images=fb_images if fb_images else None,
                        extraction_method="pptx_xml_fallback",
                        size_bytes=file_path.stat().st_size,
                        metadata={"slides": len(slide_paths), "media_images": len(fb_images)},
                    )
        except Exception as exc:
            xml_error = exc

    messages: list[str] = []
    if pptx_error:
        messages.append(f"python-pptx: {pptx_error}")
    if xml_error:
        messages.append(f"xml_fallback: {xml_error}")

    return ExtractedContent(
        filename=file_path.name,
        file_type="error",
        error="; ".join(messages) if messages else "Unsupported or corrupted PowerPoint file",
        extraction_method="failed",
        size_bytes=file_path.stat().st_size,
    )


def _detect_language(ext: str) -> str:
    """Detect programming language from extension."""
    lang_map = {
        ".py": "python",
        ".java": "java",
        ".cpp": "cpp",
        ".c": "c",
        ".h": "c",
        ".js": "javascript",
        ".ts": "typescript",
        ".jsx": "jsx",
        ".tsx": "tsx",
        ".cs": "csharp",
        ".go": "go",
        ".rb": "ruby",
        ".php": "php",
        ".swift": "swift",
        ".kt": "kotlin",
        ".scala": "scala",
        ".rs": "rust",
        ".sh": "bash",
        ".sql": "sql",
        ".html": "html",
        ".css": "css",
        ".r": "r",
        ".m": "matlab",
        ".lua": "lua",
        ".pl": "perl",
        ".dart": "dart",
        ".vue": "vue",
        ".svelte": "svelte",
    }
    return lang_map.get(ext.lower(), "unknown")


def process_student_submission(student_dir: Path, student_id: str, 
                               session_id: int) -> Tuple[List[ExtractedContent], IngestionReport]:
    """
    Process all files in a student's submission directory.
    
    Returns (list of ExtractedContent, full IngestionReport)
    """
    timestamp = datetime.now().isoformat()
    
    report = IngestionReport(
        student_id=student_id,
        session_id=session_id,
        timestamp=timestamp
    )
    
    extracted_contents = []
    
    # Find all files recursively.  Apply a targeted cross-contamination guard:
    # skip files that live inside a subdirectory whose name starts with a roll
    # number DIFFERENT from the current student's roll number.
    #
    # Examples of what we WANT to allow:
    #   Lab Task 3/file.ipynb          — subject folder, not a student
    #   22P-9171/Lab3_T1.ipynb         — student's own subfolder
    #   src/main.py                    — standard code layout
    #
    # Examples of what we WANT to block (teacher bundled all ZIPs into one ZIP,
    # so other students' folders appear as subdirs):
    #   22P-9040_Ali/work.py           — clearly a different student
    import re as _re_parser

    # Matches roll-number prefix ONLY: 22P-9171, 22p9075, 21I-0001, etc.
    _roll_pattern = _re_parser.compile(r'^\d{2}[pPiI][-_]?\d{3,5}', _re_parser.IGNORECASE)

    # Normalise the current student's roll prefix for comparison
    def _norm_roll(s: str) -> str:
        return _re_parser.sub(r'[-_]', '', s).lower()

    _current_roll = ""
    _m = _roll_pattern.match(student_dir.name)
    if _m:
        _current_roll = _norm_roll(_m.group(0))

    all_files = []
    for file_path in sorted(student_dir.rglob("*")):
        if file_path.is_file():
            # Skip hidden and system files
            rel_parts = file_path.relative_to(student_dir).parts
            if any(part.startswith(".") or part.startswith("__") for part in rel_parts):
                continue
            if file_path.name in {".DS_Store", "Thumbs.db", ".gitignore"}:
                continue
            if _is_transient_or_system_file(file_path):
                continue

            # Cross-contamination guard: skip files inside a subdirectory whose
            # roll-number prefix is DIFFERENT from the current student's roll.
            # Folders like "Lab Task 3", "src", "code" never match the roll
            # pattern so they are always allowed.
            if len(rel_parts) >= 2:
                skip = False
                for dir_part in rel_parts[:-1]:
                    dm = _roll_pattern.match(dir_part)
                    if dm:
                        dir_roll = _norm_roll(dm.group(0))
                        # Same student wrapped files in their own named subfolder — allow.
                        # Different roll number — this is a contaminated submission.
                        if _current_roll and dir_roll != _current_roll:
                            skip = True
                            break
                if skip:
                    logger.info(
                        f"[FILE_PARSER] Skipping nested student file: "
                        f"{'/'.join(rel_parts)} (probable cross-contamination)"
                    )
                    report.warnings.append(
                        f"Nested student folder detected: {rel_parts[0]} — "
                        f"skipped {'/'.join(rel_parts)} to prevent cross-contamination"
                    )
                    continue

            all_files.append(file_path)
    
    # Record received files
    for file_path in all_files:
        report.files_received.append({
            "filename": file_path.name,
            "relative_path": str(file_path.relative_to(student_dir)),
            "size_bytes": file_path.stat().st_size
        })
    
    # Parse each file
    for file_path in all_files:
        content, file_report = parse_file_with_report(file_path)
        extracted_contents.append(content)
        
        if file_report["status"] == "parsed":
            report.files_parsed.append(file_report)
            report.total_text_chars += file_report.get("text_length", 0)
            report.total_images += file_report.get("image_count", 0)
        elif file_report["status"] == "failed":
            report.files_failed.append(file_report)
            report.errors.append(f"{file_path.name}: {file_report.get('error', 'Unknown error')}")
        elif file_report.get("warning"):
            report.warnings.append(f"{file_path.name}: {file_report['warning']}")
    
    # Check for truncation
    report.content_truncated = any(
        c.metadata.get("truncated", False) for c in extracted_contents
    )
    
    return extracted_contents, report


def prepare_content_for_llm(extracted_contents: List[ExtractedContent], 
                           max_text_chars: int = 80000,
                           max_images: int = 30) -> Tuple[str, List[Dict], Dict[str, Any]]:
    """
    Prepare extracted content for LLM consumption.
    
    Returns (combined_text, list_of_images, metadata)
    """
    text_parts = []
    all_images = []
    metadata = {
        "files_processed": [],
        "truncated": False,
        "images_included": 0,
        "text_chars": 0
    }
    
    total_text_chars = 0
    
    for content in extracted_contents:
        file_info = {
            "filename": content.filename,
            "file_type": content.file_type,
            "extraction_method": content.extraction_method
        }
        
        # Add text content
        if content.text_content:
            header = f"\n=== {content.filename} ({content.file_type}) ===\n"
            text_to_add = header + content.text_content
            
            if total_text_chars + len(text_to_add) > max_text_chars:
                # Truncate
                remaining = max_text_chars - total_text_chars - len(header)
                if remaining > 100:
                    text_parts.append(header + content.text_content[:remaining])
                    text_parts.append("\n[CONTENT TRUNCATED DUE TO LENGTH LIMITS]\n")
                metadata["truncated"] = True
                break
            else:
                text_parts.append(text_to_add)
                total_text_chars += len(text_to_add)
        
        # Add images
        for img in content.images:
            if len(all_images) >= max_images:
                metadata["truncated"] = True
                break
            all_images.append({
                "filename": content.filename,
                **img
            })
        
        metadata["files_processed"].append(file_info)
    
    metadata["images_included"] = len(all_images)
    metadata["text_chars"] = total_text_chars
    
    return "\n".join(text_parts), all_images, metadata


# Backwards compatibility
parse_file = lambda file_path: parse_file_with_report(file_path)[0]
def get_file_type_summary(files: list) -> dict:
    """Get a summary of file types in a submission."""
    summary = {
        "code": [],
        "images": [],
        "pdfs": [],
        "documents": [],
        "text": [],
        "notebooks": [],
        "other": [],
        "errors": [],
        "total_count": len(files),
        "total_size": 0,
    }
    
    for f in files:
        if isinstance(f, ExtractedContent):
            file_type = f.file_type
            filename = f.filename
            size = f.size_bytes
        else:
            file_type = f.get("type", "unknown")
            filename = f.get("filename", "unknown")
            size = f.get("size", 0)
        
        summary["total_size"] += size
        
        if file_type == "code":
            summary["code"].append(filename)
        elif file_type == "image":
            summary["images"].append(filename)
        elif file_type == "pdf":
            summary["pdfs"].append(filename)
        elif file_type in ("docx", "excel", "powerpoint"):
            summary["documents"].append(filename)
        elif file_type == "text":
            summary["text"].append(filename)
        elif file_type == "notebook":
            summary["notebooks"].append(filename)
        elif file_type == "error":
            summary["errors"].append(f"{filename}: {f.error if isinstance(f, ExtractedContent) else 'unknown error'}")
        else:
            summary["other"].append(filename)
    
    return summary
